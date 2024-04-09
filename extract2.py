import json
from pptx import Presentation

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}
        actions = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:  # Check if text is not empty
                        actions.append(text)

        if actions:
            slide_data["title"] = actions[0]
            if len(actions) > 1:
                slide_data["events"] = [{"title": actions[1], "actions": actions[2:]}]
            else:
                slide_data["events"] = [{"title": None, "actions": []}]
        
        slides_data.append(slide_data)

    return slides_data

def save_to_json(data, json_file):
    with open(json_file, 'w') as f:
        json.dump(data, f, indent=4)

def main():
    ppt_file = 'type2-2.pptx'
    json_file = 'output2-2.json'

    slides_data = extract_text_from_ppt(ppt_file)
    save_to_json(slides_data, json_file)
    print("Text extracted and saved to JSON successfully.")

if __name__ == "__main__":
    main()
