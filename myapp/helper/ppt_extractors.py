from pptx import Presentation


def extract_text_from_ppt_1(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}
        slide_data['title'] = slide.shapes[0].text if slide.shapes[0].has_text_frame else None
        slide_data['events'] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes[0]:
                actions = shape.text.split('\n')
                event_title = actions.pop(0).strip()
                event_actions = [action.strip() for action in actions if action.strip()]
                slide_data['events'].append({'title': event_title, 'actions': event_actions})

        slides_data.append(slide_data)

    return slides_data

def extract_text_from_ppt_2(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}
        actions = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:  # Check if text is not empty
                        actions.append(text)

        if actions:
            slide_data["title"] = actions[0]
            if len(actions) > 1:
                slide_data["events"] = [{"title": actions[1], "actions": actions[2:]}]
            else:
                slide_data["events"] = [{"title": None, "actions": []}]
        
        slides_data.append(slide_data)

    return slides_data

# concept type
def extract_text_from_ppt_3(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}

        if slide.shapes:
            slide_text = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame]
            if slide_text:
                slide_data["name"] = slide_text[0]
                description = slide_text[1:-1] if len(slide_text) > 2 else []
                formatted_description = []
                for line in description:
                    formatted_lines = [l for l in line.split('\n') if l.strip()]  # Filter out empty lines
                    formatted_description.extend(formatted_lines[1:])  # Exclude the first index
                slide_data["description"] = formatted_description

        slides_data.append(slide_data)

    return slides_data

