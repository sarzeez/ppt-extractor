import io
from django.http import HttpResponse
from pptx import Presentation
from openpyxl import Workbook
from django.shortcuts import redirect, render

from django.conf import settings
from django.core.files.storage import FileSystemStorage

# Create your views here.
def home(request):
    if request.method == 'POST' and request.FILES:
        selected_format = request.POST.get('format')
        print(selected_format)
        uploaded_file = request.FILES['uploaded_file']
        filename = uploaded_file.name.split('.')[0]

        if selected_format == 'format1':
            data = extract_text_from_ppt(uploaded_file)
            response = HttpResponse(generate_excel(data), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            content_disposition = "attachment; filename=" + filename + ".xlsx"
            response['Content-Disposition'] = content_disposition
            return response
        elif selected_format == 'format2':
            return redirect('home')
        else:
            return redirect('home')

    return render(request, 'home.html', {})

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}
        slide_data['title'] = slide.shapes[0].text if slide.shapes[0].has_text_frame else None
        slide_data['events'] = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape != slide.shapes[0]:
                actions = shape.text.split('\n')
                event_title = actions.pop(0).strip()
                event_actions = [action.strip() for action in actions if action.strip()]
                slide_data['events'].append({'title': event_title, 'actions': event_actions})

        slides_data.append(slide_data)

    return slides_data

def generate_excel(data):
    wb = Workbook()
    ws = wb.active

    # Write headers
    ws.append(['Concept name', 'Attraction Title', 'Attraction Text'])

    # Write data
    for slide in data:
        for event in slide['events']:
            actions_with_period = [action.strip() + '.' if not action.strip().endswith('.') else action.strip() for action in event['actions']]
            actions_str = ' '.join(actions_with_period)
            ws.append([slide['title'], event['title'], actions_str])

    # Save the workbook to a BytesIO object
    excel_buffer = io.BytesIO()
    wb.save(excel_buffer)
    excel_buffer.seek(0)
    return excel_buffer
