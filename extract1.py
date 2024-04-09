import csv
import json
from pptx import Presentation

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}
        slide_data['title'] = None
        slide_data['events'] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    if slide_data['title'] is None:
                        splitted_title = text.split('\n').pop(0).strip()
                        print(splitted_title)
                        slide_data['title'] = splitted_title
                    else:
                        actions = text.split('\n')
                        event_title = actions.pop(0).strip()
                        event_actions = [action.strip() for action in actions if action.strip()]
                        slide_data['events'].append({'title': event_title, 'actions': event_actions})

        slides_data.append(slide_data)

    return slides_data


def generate_csv(data, filename):
    # Specify the file path
    file_path = filename + ".csv"

    # Open the file in write mode
    with open(file_path, 'w', newline='', encoding='utf-8') as csv_file:
        # Create a CSV writer
        csv_writer = csv.writer(csv_file)

        # Write headers
        csv_writer.writerow(['Concept Name', 'Attraction Title', 'Attraction Text'])

        # Write data
        for slide in data:
            for event in slide['events']:
                actions_str = '\n'.join(event['actions'])
                csv_writer.writerow([slide['title'], event['title'], actions_str])

    return file_path


def save_to_json(data, json_file):
    with open(json_file, 'w') as f:
        json.dump(data, f, indent=4)

def main():
    ppt_file = 'type1-cn.pptx'
    json_file = 'output1.json'

    slides_data = extract_text_from_ppt(ppt_file)
    # save_to_json(slides_data, json_file)
    generate_csv(slides_data, json_file)
    print("Text extracted and saved to JSON successfully.")

if __name__ == "__main__":
    main()
