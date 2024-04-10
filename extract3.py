import json
from pptx import Presentation

def extract_text_from_ppt(ppt_file):
    presentation = Presentation(ppt_file)
    slides_data = []

    for slide in presentation.slides:
        slide_data = {}

        if slide.shapes:
            slide_text = [shape.text.strip() for shape in slide.shapes if shape.has_text_frame]
            if slide_text:
                slide_data["name"] = slide_text[0]
                description = slide_text[1:-1] if len(slide_text) > 2 else []
                formatted_description = []
                for line in description:
                    formatted_lines = [l for l in line.split('\n') if l.strip()]  # Filter out empty lines
                    formatted_description.extend(formatted_lines[1:])  # Exclude the first index
                slide_data["description"] = formatted_description

        slides_data.append(slide_data)

    return slides_data



def save_to_json(data, json_file):
    with open(json_file, 'w') as f:
        json.dump(data, f, indent=4)

def main():
    ppt_file = 'concepts.pptx'
    json_file = 'output3.json'

    slides_data = extract_text_from_ppt(ppt_file)
    save_to_json(slides_data, json_file)
    print("Text extracted and saved to JSON successfully.")

if __name__ == "__main__":
    main()
